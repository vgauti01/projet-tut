"""Tests for document extractors."""
import sys
import tempfile
from pathlib import Path

import pytest
from pathlib import Path

from ingest.extractors.base import Extractor, ExtractedPage
from ingest.extractors.txt_extractor import TxtExtractor
from ingest.extractors.csv_extractor import CsvExtractor
from ingest.extractors.registry import get_extractor, supported_extensions


class TestExtractedPage:
    def test_dataclass(self):
        page = ExtractedPage(page_number=1, text="hello", metadata={"source_type": "test"})
        assert page.page_number == 1
        assert page.text == "hello"
        assert page.metadata["source_type"] == "test"

    def test_default_metadata(self):
        page = ExtractedPage(page_number=1, text="hello")
        assert page.metadata == {}

    def test_docling_document_default_none(self):
        page = ExtractedPage(page_number=1, text="hello")
        assert page.docling_document is None

    def test_docling_document_set(self):
        fake_doc = {"type": "docling"}
        page = ExtractedPage(page_number=1, text="hello", docling_document=fake_doc)
        assert page.docling_document == {"type": "docling"}

    def test_docling_document_not_in_repr(self):
        page = ExtractedPage(page_number=1, text="hello", docling_document="large_object")
        assert "large_object" not in repr(page)


class TestTxtExtractor:
    def test_can_handle(self):
        assert TxtExtractor.can_handle(Path("file.txt"))
        assert TxtExtractor.can_handle(Path("file.md"))
        assert TxtExtractor.can_handle(Path("file.rst"))
        assert not TxtExtractor.can_handle(Path("file.pdf"))

    def test_extract(self):
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Hello world\nSecond line")
            f.flush()
            path = Path(f.name)

        extractor = TxtExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) == 1
        assert pages[0].page_number == 1
        assert "Hello world" in pages[0].text
        assert pages[0].metadata["source_type"] == "txt"
        path.unlink()

    def test_extract_empty(self):
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("   ")
            f.flush()
            path = Path(f.name)

        extractor = TxtExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) == 0
        path.unlink()


class TestCsvExtractor:
    def test_can_handle(self):
        assert CsvExtractor.can_handle(Path("data.csv"))
        assert not CsvExtractor.can_handle(Path("data.xlsx"))

    def test_extract(self):
        with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False, encoding="utf-8") as f:
            f.write("name,age\nAlice,30\nBob,25\n")
            f.flush()
            path = Path(f.name)

        extractor = CsvExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) == 1
        assert "name: Alice" in pages[0].text
        assert "age: 30" in pages[0].text
        assert pages[0].metadata["source_type"] == "csv"
        path.unlink()

    def test_extract_batching(self, monkeypatch):
        # On réduit la taille du batch pour tester le découpage
        import ingest.extractors.csv_extractor
        monkeypatch.setattr(ingest.extractors.csv_extractor, "BATCH_SIZE", 2)

        with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False, encoding="utf-8") as f:
            f.write("id,val\n1,a\n2,b\n3,c\n4,d\n5,e\n")
            f.flush()
            path = Path(f.name)

        extractor = CsvExtractor()
        pages = list(extractor.extract(path))
        
        # 5 lignes, batch de 2 -> 3 pages (2, 2, 1)
        assert len(pages) == 3
        assert "id: 1" in pages[0].text
        assert "id: 3" in pages[1].text
        assert "id: 5" in pages[2].text
        path.unlink()


class TestDocxExtractor:
    def test_can_handle(self):
        from ingest.extractors.docx_extractor import DocxExtractor
        assert DocxExtractor.can_handle(Path("doc.docx"))
        assert not DocxExtractor.can_handle(Path("doc.txt"))

    def test_extract(self):
        from docx import Document
        from ingest.extractors.docx_extractor import DocxExtractor
        
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            doc.add_paragraph("Hello from Word")
            doc.save(f.name)
            path = Path(f.name)

        extractor = DocxExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) == 1
        assert "Hello from Word" in pages[0].text
        assert pages[0].metadata["source_type"] == "docx"
        path.unlink()


class TestXlsxExtractor:
    def test_can_handle(self):
        from ingest.extractors.xlsx_extractor import XlsxExtractor
        assert XlsxExtractor.can_handle(Path("file.xlsx"))
        assert not XlsxExtractor.can_handle(Path("file.xls"))

    def test_extract(self):
        from openpyxl import Workbook
        from ingest.extractors.xlsx_extractor import XlsxExtractor
        
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = Workbook()
            ws = wb.active
            ws.title = "TestSheet"
            ws.append(["Header1", "Header2"])
            ws.append(["Val1", "Val2"])
            wb.save(f.name)
            path = Path(f.name)

        extractor = XlsxExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) == 1
        # Docling produit du markdown tabulaire
        assert "Header1" in pages[0].text
        assert "Val1" in pages[0].text
        assert pages[0].metadata["source_type"] == "xlsx"
        path.unlink()

    def test_extract_empty_sheet(self):
        from openpyxl import Workbook
        from ingest.extractors.xlsx_extractor import XlsxExtractor
        
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = Workbook()
            ws = wb.active
            # On ne met rien dans la feuille
            wb.save(f.name)
            path = Path(f.name)

        extractor = XlsxExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) == 0
        path.unlink()


class TestPptxExtractor:
    def test_can_handle(self):
        from ingest.extractors.pptx_extractor import PptxExtractor
        assert PptxExtractor.can_handle(Path("pres.pptx"))

    def test_extract(self):
        from pptx import Presentation
        from ingest.extractors.pptx_extractor import PptxExtractor
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = "Hello PPTX"
            
            # Ajouter des notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Speaker notes"
            
            prs.save(f.name)
            path = Path(f.name)

        extractor = PptxExtractor()
        pages = list(extractor.extract(path))
        assert len(pages) >= 1
        # Docling extrait le titre, les notes peuvent ne pas être incluses
        all_text = " ".join(p.text for p in pages)
        assert "Hello PPTX" in all_text
        assert pages[0].metadata["source_type"] == "pptx"
        path.unlink()


class TestPdfExtractor:
    def test_can_handle(self):
        from ingest.extractors.pdf_extractor import PdfExtractor
        assert PdfExtractor.can_handle(Path("doc.pdf"))

    def test_extract_mock(self, monkeypatch):
        """Test with mocking because creating a valid PDF is complex without extra libs."""
        from ingest.extractors.pdf_extractor import PdfExtractor
        
        class MockPage:
            def get_text(self, opt): return "Page content"
            
        class MockDoc:
            def __iter__(self): return iter([MockPage()])
            def close(self): pass

        monkeypatch.setattr("fitz.open", lambda x: MockDoc())
        
        extractor = PdfExtractor()
        pages = list(extractor.extract(Path("fake.pdf")))
        assert len(pages) == 1
        assert pages[0].text == "Page content"
        assert pages[0].page_number == 1


class TestRegistry:
    def test_get_extractor_txt(self):
        ext = get_extractor(Path("readme.txt"))
        assert ext is not None
        assert isinstance(ext, TxtExtractor)

    def test_get_extractor_unknown(self):
        ext = get_extractor(Path("file.xyz"))
        assert ext is None

    def test_supported_extensions(self):
        exts = supported_extensions()
        assert ".pdf" in exts
        assert ".docx" in exts
        assert ".xlsx" in exts
        assert ".pptx" in exts
        assert ".csv" in exts
        assert ".txt" in exts

    def test_get_extractor_pdf(self):
        ext = get_extractor(Path("doc.pdf"))
        assert ext is not None

    def test_get_extractor_docx(self):
        ext = get_extractor(Path("doc.docx"))
        assert ext is not None
