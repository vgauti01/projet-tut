from pathlib import Path
from typing import Iterator
import logging

from .base import Extractor, ExtractedPage
from .hf_compat import patch_hf_hub

patch_hf_hub()

logger = logging.getLogger(__name__)


class PptxExtractor(Extractor):
    """
    Extracteur pour les fichiers PPTX utilisant Docling.

    Docling préserve la structure sémantique des présentations PowerPoint :
    - Détection automatique du layout et des zones de texte
    - Extraction des tableaux dans les slides avec structure préservée
    - Génération markdown structurée pour chaque slide
    - Préservation des titres et sous-titres
    - Extraction des notes du présentateur

    Fallback automatique vers python-pptx si Docling échoue.
    """
    SUPPORTED_EXTENSIONS = [".pptx"]

    def __init__(self):
        self._converter = None

    def _get_converter(self):
        """Lazy loading du DocumentConverter."""
        if self._converter is None:
            try:
                from docling.document_converter import DocumentConverter
                self._converter = DocumentConverter()
            except ImportError as e:
                logger.error("Docling n'est pas installé. Installez-le avec: uv add docling")
                raise ImportError(
                    "docling est requis pour l'extraction PPTX. "
                    "Installez-le avec: uv add docling"
                ) from e
        return self._converter

    def extract(self, file_path: Path) -> Iterator[ExtractedPage]:
        """
        Extrait le contenu du fichier PPTX en préservant la structure.

        Utilise Docling pour :
        - Analyser le layout de chaque slide
        - Détecter et extraire les tableaux avec leur structure
        - Convertir en markdown structuré par slide
        """
        try:
            converter = self._get_converter()

            # Conversion du document avec Docling
            result = converter.convert(str(file_path))

            # Extraction du texte complet avec structure préservée
            full_markdown = result.document.export_to_markdown()

            if not full_markdown or not full_markdown.strip():
                logger.warning(f"Aucun contenu extrait de {file_path}")
                return

            # Pour PPTX avec Docling, on peut essayer de détecter les pages
            # Si le document a des pages identifiées, les traiter séparément
            if hasattr(result.document, 'pages') and result.document.pages:
                for page_idx, page in enumerate(result.document.pages, start=1):
                    try:
                        # Extraction du contenu de la page
                        page_markdown = page.export_to_markdown() if hasattr(page, 'export_to_markdown') else str(page)

                        metadata = {
                            "source_type": "pptx",
                            "extraction_method": "docling",
                            "slide_number": page_idx,
                        }

                        # Ajout du titre de la slide si disponible
                        if hasattr(page, 'title') and page.title:
                            metadata["slide_title"] = page.title

                        yield ExtractedPage(
                            page_number=page_idx,
                            text=page_markdown,
                            metadata=metadata,
                        )
                    except Exception as e:
                        logger.error(f"Erreur lors de l'extraction de la slide {page_idx}: {e}")
                        continue
            else:
                # Fallback : traiter comme un document unique
                # On peut essayer de splitter par marqueurs de slide
                slides = self._split_slides_from_markdown(full_markdown)

                if slides:
                    for idx, slide_text in enumerate(slides, start=1):
                        yield ExtractedPage(
                            page_number=idx,
                            text=slide_text,
                            metadata={
                                "source_type": "pptx",
                                "extraction_method": "docling",
                            },
                        )
                else:
                    # Dernier fallback : tout comme une seule page
                    yield ExtractedPage(
                        page_number=1,
                        text=full_markdown,
                        metadata={
                            "source_type": "pptx",
                            "extraction_method": "docling",
                        },
                    )

        except Exception as e:
            logger.error(f"Impossible d'extraire le fichier PPTX {file_path} avec Docling: {e}")
            # Fallback vers python-pptx
            logger.info(f"Tentative de fallback vers python-pptx pour {file_path}")
            try:
                yield from self._fallback_python_pptx_extract(file_path)
            except Exception as fallback_error:
                logger.error(f"Le fallback python-pptx a également échoué: {fallback_error}")
                return

    def _split_slides_from_markdown(self, markdown: str) -> list[str]:
        """
        Tente de diviser le markdown en slides individuelles.
        Recherche des marqueurs de séparation de slides.
        """
        # Docling peut utiliser des marqueurs comme "---" ou des titres de niveau 1
        # pour séparer les slides
        slides = []

        # Essai 1 : Split par "---" (séparateur markdown classique)
        if "\n---\n" in markdown:
            slides = [s.strip() for s in markdown.split("\n---\n") if s.strip()]

        # Essai 2 : Split par titres de niveau 1 (# Titre)
        elif "\n# " in markdown:
            parts = markdown.split("\n# ")
            # Le premier élément peut être vide ou contenir du texte avant le premier titre
            for i, part in enumerate(parts):
                if i == 0 and not part.strip():
                    continue
                # Rajouter le # pour les titres (sauf le premier si c'était du texte initial)
                slide_text = part if i == 0 else f"# {part}"
                if slide_text.strip():
                    slides.append(slide_text.strip())

        return slides

    def _fallback_python_pptx_extract(self, file_path: Path) -> Iterator[ExtractedPage]:
        """
        Méthode de fallback utilisant python-pptx si Docling échoue.
        Conserve l'ancienne logique d'extraction.
        """
        from pptx import Presentation

        # Parse la présentation PPTX et extrait le texte de chaque slide
        prs = Presentation(str(file_path))

        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts = []

            # Extraction du texte des formes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)

            # Extraction des notes du présentateur
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Notes] {notes}")

            if parts:
                text = "\n".join(parts)
                yield ExtractedPage(
                    page_number=slide_idx,
                    text=text,
                    metadata={
                        "source_type": "pptx",
                        "extraction_method": "python_pptx_fallback",
                    },
                )
